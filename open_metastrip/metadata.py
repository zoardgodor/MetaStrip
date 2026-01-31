import os
from PIL import Image
import piexif
import mutagen
from PyPDF2 import PdfReader, PdfWriter
from datetime import datetime
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import cv2

SUPPORTED_IMAGE_FORMATS = {'.jpg', '.jpeg', '.tiff', '.tif', '.png', '.bmp', '.gif', '.webp', '.ico', '.svg'}
SUPPORTED_AUDIO_FORMATS = {'.mp3', '.flac', '.ogg', '.wav'}
SUPPORTED_PDF_FORMATS = {'.pdf'}
SUPPORTED_DOCUMENT_FORMATS = {'.docx'}
SUPPORTED_PRESENTATION_FORMATS = {'.pptx'}
SUPPORTED_SPREADSHEET_FORMATS = {'.xlsx'}
SUPPORTED_VIDEO_FORMATS = {'.mp4', '.avi', '.mkv', '.mov', '.flv', '.wmv', '.webm'}


def get_file_extension(file_path):
    return os.path.splitext(file_path)[1].lower()


def get_file_info(file_path):
    file_info = {}
    try:
        stat = os.stat(file_path)
        file_info["FILE:Fájlnév"] = os.path.basename(file_path)
        file_info["FILE:Teljes útvonal"] = file_path
        file_info["FILE:Fájlméret"] = format_file_size(stat.st_size)
        file_info["FILE:Méret (bájt)"] = str(stat.st_size)
        
        mod_time = datetime.fromtimestamp(stat.st_mtime)
        file_info["FILE:Módosítás dátuma"] = mod_time.strftime("%Y:%m:%d %H:%M:%S")
        
        access_time = datetime.fromtimestamp(stat.st_atime)
        file_info["FILE:Hozzáférés dátuma"] = access_time.strftime("%Y:%m:%d %H:%M:%S")
        
        change_time = datetime.fromtimestamp(stat.st_ctime)
        file_info["FILE:Módosítás időpontja"] = change_time.strftime("%Y:%m:%d %H:%M:%S")
        
        file_info["FILE:Fájltípus"] = detect_file_type(file_path)
        
    except Exception as e:
        file_info["FILE:Hiba"] = str(e)
    
    return file_info


def format_file_size(size_bytes):
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.2f} TB"


def detect_file_type(file_path):
    try:
        from PIL import Image
        img = Image.open(file_path)
        format_str = img.format or "Unknown"
        mime_map = {
            'JPEG': 'image/jpeg',
            'PNG': 'image/png',
            'TIFF': 'image/tiff',
            'GIF': 'image/gif'
        }
        mime = mime_map.get(format_str, 'image/' + format_str.lower())
        return f"{format_str} ({mime})"
    except:
        return "Unknown"


def extract_image_file_info(file_path):
    file_info = {}
    try:
        img = Image.open(file_path)
        file_info["FILE:Képméretek"] = f"{img.width}x{img.height}"
        file_info["FILE:Megapixel"] = f"{(img.width * img.height) / 1000000:.1f}"
        file_info["FILE:Képformátum"] = img.format or "Unknown"
        file_info["FILE:Képmód"] = img.mode
        
        if hasattr(img, 'info'):
            if 'dpi' in img.info:
                dpi = img.info['dpi']
                file_info["FILE:DPI"] = f"{dpi[0]}x{dpi[1]}"
    except Exception as e:
        file_info["FILE:Hiba"] = str(e)
    
    return file_info


def extract_metadata(file_path):
    ext = get_file_extension(file_path)
    
    meta = get_file_info(file_path)
    
    if ext in SUPPORTED_IMAGE_FORMATS:
        meta.update(extract_image_file_info(file_path))
        meta.update(extract_image_metadata(file_path))
    elif ext in SUPPORTED_AUDIO_FORMATS:
        meta.update(extract_audio_metadata(file_path))
    elif ext in SUPPORTED_PDF_FORMATS:
        meta.update(extract_pdf_metadata(file_path))
    elif ext in SUPPORTED_DOCUMENT_FORMATS:
        meta.update(extract_docx_metadata(file_path))
    elif ext in SUPPORTED_PRESENTATION_FORMATS:
        meta.update(extract_pptx_metadata(file_path))
    elif ext in SUPPORTED_SPREADSHEET_FORMATS:
        meta.update(extract_xlsx_metadata(file_path))
    elif ext in SUPPORTED_VIDEO_FORMATS:
        meta.update(extract_video_metadata(file_path))
    else:
        meta["Nem támogatott fájltípus"] = ext
    
    return meta


def extract_image_metadata(file_path):
    try:
        img = Image.open(file_path)
        exif_data = img.info.get('exif')
        if exif_data:
            exif_dict = piexif.load(exif_data)
            meta = {}
            for ifd in exif_dict:
                if exif_dict[ifd] is None:
                    continue
                for tag in exif_dict[ifd]:
                    try:
                        tag_name = piexif.TAGS[ifd][tag]["name"]
                        value = exif_dict[ifd][tag]
                        meta[f"EXIF:{tag_name}"] = value
                    except Exception:
                        continue
            return meta if meta else {"EXIF": "Nincs EXIF metaadat"}
        else:
            return {"EXIF": "Nincs EXIF metaadat"}
    except Exception as e:
        return {"Hiba": str(e)}


def extract_audio_metadata(file_path):
    try:
        audio = mutagen.File(file_path, easy=True)
        if not audio:
            return {"ID3": "Nem olvasható vagy nincs metaadat"}
        meta = {}
        items = getattr(audio, 'items', None)
        if callable(items):
            for k, v in audio.items():
                meta[f"ID3:{k}"] = v
        return meta if meta else {"ID3": "Nincs metaadat"}
    except Exception as e:
        return {"Hiba": str(e)}


def extract_pdf_metadata(file_path):
    try:
        reader = PdfReader(file_path)
        info = reader.metadata
        if info and hasattr(info, 'items'):
            meta = {f"PDF:{k}": v for k, v in info.items()}
            return meta if meta else {"PDF": "Nincs metaadat"}
        else:
            return {"PDF": "Nincs metaadat"}
    except Exception as e:
        return {"Hiba": str(e)}


def remove_metadata(file_path):
    ext = get_file_extension(file_path)
    if ext in SUPPORTED_IMAGE_FORMATS:
        return remove_image_metadata(file_path)
    elif ext in SUPPORTED_AUDIO_FORMATS:
        return remove_audio_metadata(file_path)
    elif ext in SUPPORTED_PDF_FORMATS:
        return remove_pdf_metadata(file_path)
    elif ext in SUPPORTED_DOCUMENT_FORMATS:
        return remove_docx_metadata(file_path)
    elif ext in SUPPORTED_PRESENTATION_FORMATS:
        return remove_pptx_metadata(file_path)
    elif ext in SUPPORTED_SPREADSHEET_FORMATS:
        return remove_xlsx_metadata(file_path)
    elif ext in SUPPORTED_VIDEO_FORMATS:
        return remove_video_metadata(file_path)
    else:
        return False, "Nem támogatott fájltípus"


def remove_image_metadata(file_path):
    try:
        img = Image.open(file_path)
        data = list(img.getdata())
        img_no_exif = Image.new(img.mode, img.size)
        img_no_exif.putdata(data)
        out_path = file_path  
        img_no_exif.save(out_path)
        return True, "EXIF metaadatok eltávolítva"
    except Exception as e:
        return False, str(e)


def remove_audio_metadata(file_path):
    try:
        audio = mutagen.File(file_path)
        if audio is None:
            return False, "Nem olvasható vagy nincs metaadat"
        audio.delete()
        audio.save()
        return True, "ID3/metaadatok eltávolítva"
    except Exception as e:
        return False, str(e)


def extract_docx_metadata(file_path):
    try:
        doc = Document(file_path)
        meta = {}
        
        props = doc.core_properties
        if props.title:
            meta["DOCX:Cím"] = props.title
        if props.subject:
            meta["DOCX:Tárgy"] = props.subject
        if props.author:
            meta["DOCX:Szerző"] = props.author
        if props.keywords:
            meta["DOCX:Kulcsszavak"] = props.keywords
        if props.comments:
            meta["DOCX:Megjegyzések"] = props.comments
        if props.created:
            meta["DOCX:Létrehozás dátuma"] = props.created.strftime("%Y:%m:%d %H:%M:%S")
        if props.modified:
            meta["DOCX:Módosítás dátuma"] = props.modified.strftime("%Y:%m:%d %H:%M:%S")
        
        meta["DOCX:Bekezdések száma"] = len(doc.paragraphs)
        meta["DOCX:Táblázatok száma"] = len(doc.tables)
        
        return meta if meta else {"DOCX": "Nincs metaadat"}
    except Exception as e:
        return {"Hiba": str(e)}


def extract_pptx_metadata(file_path):
    try:
        prs = Presentation(file_path)
        meta = {}
        
        props = prs.core_properties
        if props.title:
            meta["PPTX:Cím"] = props.title
        if props.subject:
            meta["PPTX:Tárgy"] = props.subject
        if props.author:
            meta["PPTX:Szerző"] = props.author
        if props.keywords:
            meta["PPTX:Kulcsszavak"] = props.keywords
        if props.comments:
            meta["PPTX:Megjegyzések"] = props.comments
        if props.created:
            meta["PPTX:Létrehozás dátuma"] = props.created.strftime("%Y:%m:%d %H:%M:%S")
        if props.modified:
            meta["PPTX:Módosítás dátuma"] = props.modified.strftime("%Y:%m:%d %H:%M:%S")
        
        meta["PPTX:Diák száma"] = len(prs.slides)
        
        return meta if meta else {"PPTX": "Nincs metaadat"}
    except Exception as e:
        return {"Hiba": str(e)}


def extract_xlsx_metadata(file_path):
    try:
        wb = load_workbook(file_path)
        meta = {}
        
        props = wb.properties
        if props.title:
            meta["XLSX:Cím"] = props.title
        if props.subject:
            meta["XLSX:Tárgy"] = props.subject
        if props.author:
            meta["XLSX:Szerző"] = props.author
        if props.keywords:
            meta["XLSX:Kulcsszavak"] = props.keywords
        if props.comments:
            meta["XLSX:Megjegyzések"] = props.comments
        if props.created:
            meta["XLSX:Létrehozás dátuma"] = props.created.strftime("%Y:%m:%d %H:%M:%S")
        if props.modified:
            meta["XLSX:Módosítás dátuma"] = props.modified.strftime("%Y:%m:%d %H:%M:%S")
        
        meta["XLSX:Munkalapok száma"] = len(wb.sheetnames)
        meta["XLSX:Munkalapok"] = ", ".join(wb.sheetnames)
        
        return meta if meta else {"XLSX": "Nincs metaadat"}
    except Exception as e:
        return {"Hiba": str(e)}


def extract_video_metadata(file_path):
    try:
        cap = cv2.VideoCapture(file_path)
        if not cap.isOpened():
            return {"VIDEO": "Videó nem olvasható"}
        
        meta = {}
        
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = cap.get(cv2.CAP_PROP_FPS)
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        
        if fps > 0:
            duration_seconds = frame_count / fps
            hours = int(duration_seconds // 3600)
            minutes = int((duration_seconds % 3600) // 60)
            seconds = int(duration_seconds % 60)
            meta["VIDEO:Hossz"] = f"{hours:02d}:{minutes:02d}:{seconds:02d}"
        
        meta["VIDEO:Képkocka szám"] = str(frame_count)
        meta["VIDEO:FPS"] = f"{fps:.2f}"
        meta["VIDEO:Felbontás"] = f"{width}x{height}"
        meta["VIDEO:Megapixel/frame"] = f"{(width * height) / 1000000:.1f}"
        meta["VIDEO:Képarány"] = f"{width}:{height}"
        
        cap.release()
        
        return meta if meta else {"VIDEO": "Nincs videó adat"}
    except Exception as e:
        return {"Hiba": str(e)}


def remove_pdf_metadata(file_path):
    try:
        reader = PdfReader(file_path)
        writer = PdfWriter()
        
        for page_num in range(len(reader.pages)):
            writer.add_page(reader.pages[page_num])
        
        writer.add_metadata({})
        
        with open(file_path, 'wb') as output_file:
            writer.write(output_file)
        
        return True, "PDF metaadatok eltávolítva"
    except Exception as e:
        return False, str(e)


def remove_docx_metadata(file_path):
    try:
        doc = Document(file_path)
        
        props = doc.core_properties
        props.title = None
        props.subject = None
        props.author = "Unknown"
        props.keywords = None
        props.comments = None
        
        doc.save(file_path)
        return True, "DOCX metaadatok eltávolítva"
    except Exception as e:
        return False, str(e)


def remove_pptx_metadata(file_path):
    try:
        prs = Presentation(file_path)
        
        props = prs.core_properties
        props.title = None
        props.subject = None
        props.author = "Unknown"
        props.keywords = None
        props.comments = None
        
        prs.save(file_path)
        return True, "PPTX metaadatok eltávolítva"
    except Exception as e:
        return False, str(e)


def remove_xlsx_metadata(file_path):
    try:
        wb = load_workbook(file_path)
        
        props = wb.properties
        props.title = None
        props.subject = None
        props.author = "Unknown"
        props.keywords = None
        props.comments = None
        
        wb.save(file_path)
        return True, "XLSX metaadatok eltávolítva"
    except Exception as e:
        return False, str(e)


def remove_video_metadata(file_path):
    return False, "Cannot remove metadata from this file type"

